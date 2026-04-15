"""
Generates PRESENTATION.pptx from Nexus project data.
Run: python generate_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ─── Palette ───────────────────────────────────────────────────────────────
SF_DARK   = RGBColor(0x03, 0x2d, 0x60)   # Salesforce navy
SF_BLUE   = RGBColor(0x01, 0x76, 0xd3)   # Salesforce blue
SF_ORANGE = RGBColor(0xfe, 0x93, 0x39)   # Salesforce orange
SF_GREEN  = RGBColor(0x2e, 0x84, 0x4a)   # Salesforce green
SF_PURPLE = RGBColor(0x75, 0x26, 0xe3)   # Salesforce purple
SF_RED    = RGBColor(0xea, 0x00, 0x1e)   # Salesforce red
WHITE     = RGBColor(0xff, 0xff, 0xff)
LIGHT_BG  = RGBColor(0xf0, 0xf4, 0xf8)
MUTED     = RGBColor(0x44, 0x44, 0x44)
BORDER    = RGBColor(0xd8, 0xdd, 0xe6)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ─── Helpers ───────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=SF_DARK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_para(tf, text, size=13, bold=False, color=MUTED,
             align=PP_ALIGN.LEFT, space_before=Pt(4), bullet=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    if bullet:
        p.level = 1
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def add_bullet_box(slide, items, x, y, w, h,
                   title=None, title_color=SF_BLUE, accent=SF_BLUE):
    # card background
    add_rect(slide, x, y, w, h, fill=WHITE, line=BORDER, line_w=Pt(1))
    # left accent bar
    add_rect(slide, x, y, 0.06, h, fill=accent)
    cur_y = y + 0.15
    if title:
        add_text(slide, title, x+0.15, cur_y, w-0.2, 0.32,
                 size=13, bold=True, color=title_color)
        cur_y += 0.35
    for item in items:
        add_text(slide, "•  " + item, x+0.15, cur_y, w-0.25, 0.28,
                 size=11, color=MUTED)
        cur_y += 0.28

def slide_header(slide, num, title, subtitle=""):
    # top bar
    add_rect(slide, 0, 0, 13.33, 0.9, fill=SF_DARK)
    # slide number circle
    circ = slide.shapes.add_shape(9, Inches(0.18), Inches(0.18),
                                   Inches(0.55), Inches(0.55))
    circ.fill.solid(); circ.fill.fore_color.rgb = SF_BLUE
    circ.line.fill.background()
    tf = circ.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = str(num)
    run.font.size = Pt(14); run.font.bold = True
    run.font.color.rgb = WHITE
    add_text(slide, title, 0.88, 0.1, 10, 0.45,
             size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.88, 0.52, 10, 0.3,
                 size=12, color=RGBColor(0xb0,0xc8,0xe8))
    # bottom accent line
    add_rect(slide, 0, 7.38, 13.33, 0.12, fill=SF_BLUE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_DARK)
# gradient-like overlay
add_rect(s, 0, 0, 6.67, 7.5, fill=RGBColor(0x0b, 0x53, 0x94))

# NEXUS
tb = s.shapes.add_textbox(Inches(1.5), Inches(1.4), Inches(10), Inches(1.6))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r1 = p.add_run(); r1.text = "NEX"
r1.font.size = Pt(80); r1.font.bold = True; r1.font.color.rgb = WHITE
r2 = p.add_run(); r2.text = "US"
r2.font.size = Pt(80); r2.font.bold = True; r2.font.color.rgb = SF_ORANGE

add_text(s, "B2B E-Commerce & CRM Platform — Built on Salesforce",
         1, 3.1, 11.33, 0.6, size=20, color=RGBColor(0xb0,0xc8,0xe8),
         align=PP_ALIGN.CENTER)

# badges row
badges = ["⚡ Salesforce Experience Cloud", "🤖 Einstein AI + Agentforce",
          "🖊 DocuSign + Dvelop E-Sign", "💳 Stripe Payments", "🇫🇷 French UI"]
bx = 0.5
for b in badges:
    bw = 2.4
    add_rect(s, bx, 4.1, bw, 0.45,
             fill=RGBColor(0x05,0x40,0x80), line=RGBColor(0x30,0x70,0xb0), line_w=Pt(1))
    add_text(s, b, bx+0.05, 4.15, bw-0.1, 0.35, size=10.5,
             color=WHITE, align=PP_ALIGN.CENTER)
    bx += bw + 0.12

add_text(s, "March 2026", 0, 6.9, 13.33, 0.4, size=11,
         color=RGBColor(0x60,0x80,0xa0), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IS SALESFORCE
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 3, "What is Salesforce?", "The world's #1 CRM platform — and much more")
add_rect(s, 0, 0.9, 13.33, 6.48, fill=LIGHT_BG)

# left column — description + table
add_rect(s, 0.3, 1.05, 6.1, 1.0, fill=SF_DARK)
add_text(s, "Salesforce is a cloud-based SaaS platform used by 150,000+ companies worldwide. "
            "Founded in 1999 — pioneered the 'No Software' movement. Everything over the internet, "
            "zero local installation.",
         0.35, 1.08, 6.0, 0.92, size=11, color=WHITE)

add_text(s, "Clouds Used in Nexus", 0.3, 2.18, 3.5, 0.3, size=13, bold=True, color=SF_DARK)
clouds = [
    ("Sales Cloud",      "Leads, Opportunities, Accounts — CRM backbone"),
    ("Experience Cloud", "Customer-facing B2B/B2C portal"),
    ("Service Cloud",    "Support cases & customer service"),
    ("Einstein AI",      "Lead scoring, recommendations, predictions"),
    ("Agentforce",       "AI agents for autonomous automation"),
]
cy = 2.55
for cloud, desc in clouds:
    add_rect(s, 0.3, cy, 6.1, 0.36, fill=WHITE, line=BORDER, line_w=Pt(1))
    add_text(s, cloud, 0.4, cy+0.04, 1.8, 0.28, size=11, bold=True, color=SF_BLUE)
    add_text(s, desc,  2.25, cy+0.04, 4.1, 0.28, size=11, color=MUTED)
    cy += 0.38

# right column — market share + concepts
add_text(s, "CRM Market Share 2024", 6.7, 1.05, 6.0, 0.3, size=13, bold=True, color=SF_DARK)
bars = [("Salesforce", 23.8, SF_BLUE), ("SAP CRM", 5.9, MUTED),
        ("Oracle", 5.3, MUTED), ("Microsoft", 4.7, MUTED), ("Adobe", 4.2, MUTED)]
by = 1.45
for name, pct, col in bars:
    add_text(s, f"{name}  {pct}%", 6.7, by, 2.2, 0.25, size=10, color=MUTED)
    bar_w = (pct / 23.8) * 3.5
    add_rect(s, 8.95, by+0.03, bar_w, 0.18, fill=col)
    by += 0.3

add_text(s, "Key Concepts", 6.7, 3.4, 6.0, 0.3, size=13, bold=True, color=SF_DARK)
concepts = ["Apex — Server-side Java-like language (backend)",
            "LWC — Modern JS/HTML component framework (frontend)",
            "SOQL — Salesforce Object Query Language (like SQL)",
            "Objects — Database tables (standard + custom)",
            "Flows — Drag-and-drop no-code automation",
            "Platform Events — Real-time event messaging",
            "Profiles & Permission Sets — Role-based access control"]
cy2 = 3.78
for c in concepts:
    add_text(s, "•  " + c, 6.7, cy2, 6.3, 0.27, size=11, color=MUTED)
    cy2 += 0.27

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — TECH STACK
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 4, "Frameworks & Technologies", "The full Nexus tech stack")
add_rect(s, 0, 0.9, 13.33, 6.48, fill=LIGHT_BG)

cards = [
    ("⚡ LWC (Frontend)", SF_BLUE,
     ["Salesforce's modern JS/HTML framework",
      ".html template + .js controller + .css",
      "@wire decorator — reactive data binding",
      "@api/@track — public API & reactivity",
      "CustomEvent — inter-component messaging",
      "NavigationMixin — page routing"]),
    ("☕ Apex (Backend)", SF_PURPLE,
     ["Server-side Java-like language",
      "Runs natively on Salesforce servers",
      "@AuraEnabled — exposes to LWC",
      "with sharing — security enforcement",
      "SOQL/DML — database operations",
      "HTTP callouts — external API calls"]),
    ("🔄 Flows (Automation)", SF_GREEN,
     ["No-code process automation builder",
      "Record-Triggered → lead AI scoring",
      "Platform Event → low-stock response",
      "Scheduled → stock risk recalculation",
      "Screen Flow → portal user setup",
      "9 flows total in Nexus"]),
]
cx = 0.3
for title, color, items in cards:
    add_rect(s, cx, 1.05, 4.1, 3.6, fill=WHITE, line=BORDER, line_w=Pt(1))
    add_rect(s, cx, 1.05, 4.1, 0.55, fill=color)
    add_text(s, title, cx+0.12, 1.1, 3.9, 0.42, size=14, bold=True, color=WHITE)
    iy = 1.7
    for item in items:
        add_text(s, "•  " + item, cx+0.15, iy, 3.8, 0.28, size=11, color=MUTED)
        iy += 0.28
    cx += 4.28

# integrations row
add_text(s, "External Integrations", 0.3, 4.85, 6, 0.3, size=13, bold=True, color=SF_DARK)
integrations = [
    ("🖊 DocuSign + Dvelop", "E-Signature (dual provider)\nREST API + Webhooks"),
    ("💳 Stripe", "Hosted Checkout\nPayment processing"),
    ("🤖 Agentforce", "AI Agents\nEinstein ML Platform"),
    ("🔗 External Negotiation", "Third-party negotiation\nREST API integration"),
]
ix = 0.3
for icon_title, desc in integrations:
    add_rect(s, ix, 5.22, 3.0, 1.05, fill=WHITE, line=BORDER, line_w=Pt(1))
    add_text(s, icon_title, ix+0.12, 5.28, 2.8, 0.32, size=12, bold=True, color=SF_DARK)
    add_text(s, desc, ix+0.12, 5.62, 2.8, 0.55, size=10, color=MUTED)
    ix += 3.18

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 5, "The Problem", "What B2B tech distributors face daily")
add_rect(s, 0, 0.9, 13.33, 6.48, fill=LIGHT_BG)

problems = [
    ("❌ Fragmented Sales Processes",
     "Quotes built in Word, negotiated via email, signed physically. No version control, no audit trail."),
    ("❌ No Real-Time Inventory Visibility",
     "Sales reps quote out-of-stock products. Nothing links quotes to actual inventory levels."),
    ("❌ No Customer Self-Service",
     "Every update requires a call or email: 'What's my quote?', 'Where's my order?', 'Can I get my invoice?'"),
    ("❌ No Data Intelligence",
     "Zero live dashboard. Pipeline health, win rates, revenue trends — all managed on gut feeling."),
    ("❌ Slow Manual Lead Qualification",
     "Without scoring, every lead gets equal attention. Reps waste hours on 5% conversion leads."),
    ("❌ No Integrated Payment Flow",
     "Checkout requires manual intervention. No online payment, no auto order creation, no invoice."),
]
px, py, pw, ph = 0.3, 1.1, 6.1, 0.88
col = 0
for title, desc in problems:
    add_rect(s, px, py, pw, ph, fill=RGBColor(0xff,0xf5,0xf5), line=RGBColor(0xff,0xd0,0xd0), line_w=Pt(1))
    add_rect(s, px, py, 0.07, ph, fill=SF_RED)
    add_text(s, title, px+0.15, py+0.05, pw-0.2, 0.3, size=12, bold=True, color=RGBColor(0xcc,0x00,0x00))
    add_text(s, desc,  px+0.15, py+0.37, pw-0.2, 0.48, size=10.5, color=RGBColor(0x66,0x00,0x00))
    col += 1
    if col % 2 == 0:
        px = 0.3; py += ph + 0.12
    else:
        px = 6.7

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — OUR SOLUTION
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 6, "Our Solution — Nexus", "One unified platform. Every problem solved.")
add_rect(s, 0, 0.9, 13.33, 6.48, fill=LIGHT_BG)

# highlight quote
add_rect(s, 0.3, 1.05, 12.73, 0.72, fill=SF_DARK)
add_text(s, "Everything runs inside Salesforce — zero external backend, zero middleware, one unified database.",
         0.45, 1.1, 12.5, 0.62, size=13, color=WHITE, align=PP_ALIGN.CENTER)

# before / after
add_rect(s, 0.3, 1.88, 5.9, 3.8, fill=RGBColor(0xff,0xf5,0xe0), line=BORDER, line_w=Pt(1))
add_rect(s, 6.5, 1.88, 6.53, 3.8, fill=RGBColor(0xe8,0xf7,0xed), line=BORDER, line_w=Pt(1))
add_text(s, "✗  Before Nexus", 0.4, 1.93, 5.7, 0.32, size=13, bold=True, color=SF_RED)
add_text(s, "✓  After Nexus",  6.6, 1.93, 6.3, 0.32, size=13, bold=True, color=SF_GREEN)

before = ["Word/Excel quote management","Email negotiation threads",
          "Physical contract signatures","No stock visibility",
          "No customer portal","Zero analytics","Manual lead qualification","Manual invoicing"]
after  = ["Versioned, AI-scored quote builder","In-portal quote messaging",
          "Automated DocuSign / Dvelop e-sign","Real-time Stock War Room",
          "Full self-service customer portal","Live KPI dashboard + pipeline funnel",
          "AI-scored leads on creation (0–100)","Auto-generated PDF invoices"]
by2 = 2.32
for b, a in zip(before, after):
    add_text(s, "•  " + b, 0.4, by2, 5.7, 0.28, size=10.5, color=MUTED)
    add_text(s, "•  " + a, 6.6, by2, 6.3, 0.28, size=10.5, color=SF_GREEN)
    by2 += 0.28

# KPI row
kpis = [("60+", "LWC Components"), ("30+", "Apex Classes"),
        ("9", "Automation Flows"), ("6", "Custom Objects"), ("4", "Integrations")]
kx = 0.3
for num, lbl in kpis:
    add_rect(s, kx, 5.82, 2.4, 0.9, fill=WHITE, line=BORDER, line_w=Pt(1))
    add_text(s, num, kx, 5.88, 2.4, 0.48, size=28, bold=True, color=SF_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, lbl, kx, 6.36, 2.4, 0.28, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    kx += 2.57

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — FUNCTIONAL REQUIREMENTS
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 7, "Functional Requirements", "What the system must do")
add_rect(s, 0, 0.9, 13.33, 6.48, fill=LIGHT_BG)

add_bullet_box(s, [
    "FR-C01 — Product catalog: filter, search, sort, grid/list view",
    "FR-C02 — Product details: specs, stock badges, ratings",
    "FR-C03 — Combo builder: slot-based hardware configurations",
    "FR-C04 — Cart & checkout: sessionStorage cart + Stripe redirect",
    "FR-C05 — Quote portal: view, accept, reject, message, download PDF",
    "FR-C06 — B2B quote request: submit form → email sales team",
    "FR-C07 — Order tracking: visual timeline, invoices, line items",
    "FR-C08 — AI recommendations: upsell/cross-sell/bundle suggestions",
    "FR-C09 — In-app notifications: bell icon, read/unread status",
    "FR-C10 — Authentication: portal login + password change",
], 0.3, 1.05, 6.2, 5.9, title="Customer Portal (FR-C)", accent=SF_BLUE, title_color=SF_BLUE)

add_bullet_box(s, [
    "FR-A01 — Quote builder: product search, stock check, AI score, send",
    "FR-A02 — Analytics: revenue, pipeline, win rate, trends, top items",
    "FR-A03 — Stock War Room: risk classification, AI summaries, PO links",
    "FR-A04 — Lead management: AI scoring, approve/reject, convert to opp",
    "FR-A05 — Order board: Kanban, status updates, auto date stamps",
    "FR-A06 — E-signature: send DocuSign/Dvelop, webhook, signed PDF",
    "FR-A07 — Quote versioning: revise without losing history",
    "FR-A08 — Convert to Order: accepted quote → order + 'Closed Won'",
    "FR-A09 — Portal user creation: account + welcome email on conversion",
    "FR-A10 — Invoice generation: auto-generated on shipment (Flow)",
], 6.83, 1.05, 6.2, 5.9, title="Admin / Sales Rep (FR-A)", accent=SF_ORANGE, title_color=SF_ORANGE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — NON-FUNCTIONAL REQUIREMENTS
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 8, "Non-Functional Requirements", "How the system must behave")
add_rect(s, 0, 0.9, 13.33, 6.48, fill=LIGHT_BG)

nfr_cards = [
    ("⚡ Performance", SF_BLUE,
     ["LWC loads within 2 seconds","Max 100 SOQL per transaction","Apex: 10-sec timeout max","10,000+ products supported","Paginated (50 per page)"]),
    ("🔐 Security", SF_RED,
     ["All Apex: with sharing","SOQL: WITH SECURITY_ENFORCED","API keys in Custom MDT (never hardcoded)","Portal: auth enforced on all data","Stripe handles PCI (no card data in SF)"]),
    ("✅ Reliability", SF_GREEN,
     ["99.9% uptime (Salesforce SLA)","DML wrapped in try-catch","External callouts: timeout + error handling","Platform Events ensure alert delivery","Flows have fault paths"]),
    ("📈 Scalability", SF_PURPLE,
     ["Multi-company data model","Bulkified Apex (no SOQL in loops)","Unlimited quote versions","Enterprise expansion ready"]),
    ("🎨 Usability", SF_ORANGE,
     ["French UI for customer portal","Inline form validation","Loading spinners on async calls","AI reasons in plain language","Desktop-optimized (1280px+)"]),
    ("🔧 Maintainability", MUTED,
     ["Git + SFDX source control","External configs in Custom MDT","Controller/Service Apex pattern","Single-responsibility LWC","GDPR-compliant via Salesforce"]),
]
cx2, cy2 = 0.3, 1.05
col2 = 0
for title, col, items in nfr_cards:
    add_rect(s, cx2, cy2, 4.1, 2.95, fill=WHITE, line=BORDER, line_w=Pt(1))
    add_rect(s, cx2, cy2, 4.1, 0.45, fill=col)
    add_text(s, title, cx2+0.12, cy2+0.07, 3.9, 0.35, size=13, bold=True, color=WHITE)
    iy2 = cy2 + 0.55
    for item in items:
        add_text(s, "•  " + item, cx2+0.15, iy2, 3.8, 0.26, size=10.5, color=MUTED)
        iy2 += 0.26
    col2 += 1
    if col2 % 3 == 0:
        cx2 = 0.3; cy2 += 3.1
    else:
        cx2 += 4.37

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CORE FEATURES
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 9, "Core Features", "The 8 pillars of Nexus")
add_rect(s, 0, 0.9, 13.33, 6.48, fill=LIGHT_BG)

features = [
    ("📋 B2B Quote Lifecycle", SF_BLUE,
     "Full cycle: creation → AI scoring → send → negotiate → accept → e-sign → order → invoice"),
    ("🤖 AI Lead Scoring", SF_PURPLE,
     "Every new lead auto-scored 0–100 on creation. Score + reasoning displayed on Lead record."),
    ("📦 Stock War Room", SF_RED,
     "Real-time inventory risk: Critical/At Risk/Healthy/Overstock. AI summaries. Auto purchase orders."),
    ("📊 Analytics Dashboard", SF_GREEN,
     "Revenue, pipeline, win rate, top products, top accounts, monthly trends, stock health — all live."),
    ("🛒 Customer Portal", SF_ORANGE,
     "Catalog, combo builder, cart, Stripe checkout, order tracking, quote portal, invoice download."),
    ("🖊 Dual E-Signature", SF_BLUE,
     "DocuSign + Dvelop integration. Webhooks handle completion. Signed PDF downloadable from portal."),
    ("🔔 Notification System", SF_PURPLE,
     "In-app bell notifications (Info/Success/Warning/Error). Quote, order, and stock alert events."),
    ("🧩 Combo Builder", SF_GREEN,
     "Slot-based hardware configurator. Smart category filtering per slot. Adds full config to cart."),
]
fx, fy = 0.3, 1.05
fcol = 0
for title, color, desc in features:
    add_rect(s, fx, fy, 3.13, 1.38, fill=WHITE, line=BORDER, line_w=Pt(1))
    add_rect(s, fx, fy, 3.13, 0.42, fill=color)
    add_text(s, title, fx+0.1, fy+0.06, 2.95, 0.3, size=12, bold=True, color=WHITE)
    add_text(s, desc, fx+0.1, fy+0.5, 2.95, 0.82, size=10, color=MUTED)
    fcol += 1
    if fcol % 4 == 0:
        fx = 0.3; fy += 1.5
    else:
        fx += 3.27

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CLASS DIAGRAM SALES SIDE
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 10, "Class Diagram — Sales Side", "Apex controllers, services and their relationships")
add_rect(s, 0, 0.9, 13.33, 6.48, fill=LIGHT_BG)

def class_box(slide, x, y, w, title, methods, color=SF_BLUE):
    h = 0.42 + len(methods) * 0.26 + 0.1
    add_rect(slide, x, y, w, h, fill=WHITE, line=color, line_w=Pt(2))
    add_rect(slide, x, y, w, 0.38, fill=color)
    add_text(slide, title, x+0.08, y+0.05, w-0.12, 0.3, size=11, bold=True, color=WHITE)
    my = y + 0.44
    for m in methods:
        add_text(slide, "+ " + m, x+0.08, my, w-0.1, 0.24, size=9, color=MUTED)
        my += 0.24
    return h

class_box(s, 0.3, 1.1, 4.0, "NexusQuoteController",
    ["getProducts(search): List<Product2>",
     "createQuote(data): TALEXO_Quote__c",
     "getAllQuotes(): List<Quote>",
     "sendToCustomer(quoteId): void",
     "createVersionFromEdits(): Quote",
     "convertToOrder(quoteId): Id",
     "generateAndSendContract(): void"])

# arrow
add_rect(s, 2.15, 3.18, 0.04, 0.35, fill=SF_BLUE)
add_text(s, "▼ calls", 1.8, 3.5, 1.0, 0.25, size=9, color=SF_BLUE, align=PP_ALIGN.CENTER)

class_box(s, 0.3, 3.72, 4.0, "QuoteGenerationService",
    ["buildQuoteRecord(data): Quote",
     "calculateAIScore(items): Number",
     "createLineItems(quoteId): void",
     "versionQuote(quoteId): Quote",
     "createOrderFromQuote(id): Order"])

class_box(s, 4.7, 1.1, 3.9, "StockIntelligenceController",
    ["getStockRiskProducts(): List<Stock>",
     "getQuotesRelatedToStock(): List",
     "getStockForProducts(ids): Map"])

add_rect(s, 6.6, 2.46, 0.04, 0.4, fill=SF_BLUE)
add_text(s, "▼ calls", 6.2, 2.85, 1.0, 0.25, size=9, color=SF_BLUE, align=PP_ALIGN.CENTER)

class_box(s, 4.7, 3.2, 3.9, "StockVelocityEngine",
    ["calculateVelocity(stockId): Decimal",
     "classifyRisk(v, qty): String",
     "generateRiskSummary(s): String"])

class_box(s, 4.7, 4.72, 3.9, "AnalyticsDashboardController",
    ["getKpiSummary(): Map",
     "getRevenueByMonth(): List",
     "getQuotePipeline(): List",
     "getTopProducts(): List",
     "getTopAccounts(): List"])

class_box(s, 8.9, 1.1, 4.13, "DvelopSignService",
    ["sendForSignature(quoteId): String",
     "discoverTenant(): String",
     "uploadDocument(pdf): String",
     "addSignatories(env, emails): void",
     "getSignedDocumentBase64(): String"])

class_box(s, 8.9, 3.2, 4.13, "LeadScoringWriter  (@InvocableMethod)",
    ["scoreLead(request): ScoreResult",
     "calculateScore(lead): Integer",
     "buildReason(lead): String"])

class_box(s, 8.9, 4.58, 4.13, "LeadManagementController",
    ["getLeadsToConvert(): List",
     "convertLeadFull(): Map",
     "createPortalUser(contactId): Id",
     "approveRejectLead(): void"])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CLASS DIAGRAM CUSTOMER SIDE
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 11, "Class Diagram — Customer Side", "Data objects and portal-facing controllers")
add_rect(s, 0, 0.9, 13.33, 6.48, fill=LIGHT_BG)

class_box(s, 0.3, 1.05, 3.9, "TALEXO_Quote__c  (Object)",
    ["Name: AutoNumber (DEVIS-0001)",
     "Opportunity__c: Lookup(Opportunity)",
     "Status__c: Picklist",
     "GrandTotal__c: Currency",
     "AcceptanceProbability__c: Number",
     "AIReason__c: LongText",
     "VersionNumber__c: Number",
     "IsPrimary__c: Boolean",
     "DvelopEnvelopeId__c: Text"],
    color=SF_PURPLE)

class_box(s, 0.3, 4.1, 3.9, "NexusQuoteLineItem__c  (Object)",
    ["Quote__c: Lookup(TALEXO_Quote__c)",
     "Product__c: Lookup(Product2)",
     "Quantity__c: Number",
     "UnitPrice__c: Currency",
     "LineTotal__c: Currency"],
    color=SF_PURPLE)

class_box(s, 4.6, 1.05, 3.9, "Product2  (Object)",
    ["Name, ProductCode, IsActive",
     "UnitPrice: Currency",
     "Subcategory__c: Text",
     "Brand__c: Text",
     "Specs__c: LongText (JSON)",
     "Rating__c: Number",
     "ImageUrl__c: Text"],
    color=SF_GREEN)

class_box(s, 4.6, 3.62, 3.9, "Stock__c  (Object)",
    ["Product__c: Lookup(Product2)",
     "Quantity_Available__c: Number",
     "Daily_Velocity__c: Number",
     "Days_Remaining__c: Number",
     "Risk_Level__c: Picklist",
     "AI_Risk_Summary__c: LongText"],
    color=SF_GREEN)

class_box(s, 8.9, 1.05, 4.13, "QuoteController  (Portal)",
    ["getAccountQuotes(accountId): List",
     "processQuoteDecision(id, d): void",
     "requestQuote(formData): void",
     "getSignedContractBase64(id): String"])

class_box(s, 8.9, 2.7, 4.13, "Order  (Object)",
    ["AccountId: Lookup(Account)",
     "Status: Picklist",
     "TotalAmount: Currency",
     "Shipping_Amount__c: Currency",
     "Shipped_Date__c: DateTime",
     "Related_Quote__c: Lookup"],
    color=SF_ORANGE)

class_box(s, 8.9, 4.78, 4.13, "AI_Recommendation__c  (Object)",
    ["Customer__c: Lookup(Account)",
     "Product__c: Lookup(Product2)",
     "Type__c: Picklist (Upsell/Cross/Bundle)",
     "Score__c: Number (0-100)",
     "Rank__c: Number"],
    color=RGBColor(0x0b, 0x79, 0x9e))

# relationship labels
add_text(s, "1──*", 2.1, 3.88, 0.8, 0.25, size=10, bold=True, color=SF_PURPLE)
add_text(s, "1──1", 4.52, 2.8, 0.7, 0.25, size=10, bold=True, color=SF_GREEN)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SEQUENCE: QUOTE LIFECYCLE
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 12, "Sequence Diagram — Quote Lifecycle", "From product selection to order + invoice")
add_rect(s, 0, 0.9, 13.33, 6.48, fill=RGBColor(0x1e,0x1e,0x2e))

# flow steps banner
steps = ["1. Build Quote","2. AI Scores","3. Send to Customer","4. Customer Accepts","5. E-Sign (Dvelop)","6. Convert → Order"]
sx = 0.2
for step in steps:
    add_rect(s, sx, 1.05, 2.05, 0.38, fill=SF_BLUE)
    add_text(s, step, sx+0.05, 1.1, 1.95, 0.28, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    sx += 2.18

# actors
actors = ["Admin", "QuoteBuilder\nLWC", "NexusQuote\nCtrl", "QuoteGen\nService", "Customer\nLWC", "Dvelop\nAPI"]
ax = [0.3, 2.0, 3.9, 5.7, 8.0, 10.3]
for i, (actor, x) in enumerate(zip(actors, ax)):
    add_rect(s, x, 1.58, 1.4, 0.48, fill=SF_DARK, line=SF_BLUE, line_w=Pt(1))
    add_text(s, actor, x, 1.6, 1.4, 0.44, size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # lifeline
    add_rect(s, x+0.68, 2.08, 0.02, 4.9, fill=RGBColor(0x30,0x50,0x80))

# sequence steps text
seq_lines = [
    (0.3, 2.18, "Admin searches products, checks stock, sets discount & terms"),
    (0.3, 2.52, "─────────────────────────────────────────────────────────▶   createQuote(data)"),
    (0.3, 2.78, "                              ──────────────────────────▶   buildQuoteRecord()"),
    (0.3, 3.04, "                                         ───────────────▶   INSERT Quote + LineItems"),
    (0.3, 3.30, "                                         ───────────────▶   calculateAIScore() → 72%"),
    (0.3, 3.56, "◀─────────────────────────────────────────────────────────   quoteId + AcceptProb"),
    (0.3, 3.82, "─────────────────────────────────────────────────────────▶   sendToCustomer()"),
    (0.3, 4.08, "                                                             UPDATE Status='Sent' + Notification"),
    (0.3, 4.34, "                                               ◀──────────   Customer sees quote, clicks Accept"),
    (0.3, 4.60, "                              ◀────────────────────────────   processQuoteDecision('Accepted')"),
    (0.3, 4.86, "─────────────────────────────────────────────────────────────────────────────────▶  POST /envelopes"),
    (0.3, 5.12, "◀─────────────────────────────────────────────────────────────────────────────────  envelopeId stored"),
    (0.3, 5.38, "─────────────────────────────────────────────────────────▶   convertToOrder() → INSERT Order"),
    (0.3, 5.64, "                                                             Opportunity → 'Closed Won'"),
]
for (lx, ly, txt) in seq_lines:
    add_text(s, txt, lx, ly, 12.9, 0.27, size=8.5, color=RGBColor(0xcd,0xd6,0xf4))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SEQUENCE: LEAD SCORING
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 13, "Sequence Diagram — AI Lead Scoring", "From form submission to portal user creation")
add_rect(s, 0, 0.9, 13.33, 6.48, fill=RGBColor(0x1e,0x1e,0x2e))

steps2 = ["1. Lead Submitted","2. Flow Triggers","3. AI Scores (0-100)","4. Admin Reviews","5. Approve & Convert","6. Portal User Created"]
sx2 = 0.2
for step in steps2:
    add_rect(s, sx2, 1.05, 2.05, 0.38, fill=SF_PURPLE)
    add_text(s, step, sx2+0.05, 1.1, 1.95, 0.28, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    sx2 += 2.18

actors2 = ["Prospect", "LeadForm\nLWC", "Lead\nController", "SF Flow", "LeadScoring\nWriter", "Admin", "LeadMgmt\nCtrl"]
ax2 = [0.15, 1.6, 3.05, 4.6, 5.9, 7.6, 9.3]
for actor, x in zip(actors2, ax2):
    add_rect(s, x, 1.58, 1.25, 0.48, fill=SF_DARK, line=SF_PURPLE, line_w=Pt(1))
    add_text(s, actor, x, 1.6, 1.25, 0.44, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x+0.61, 2.08, 0.02, 4.9, fill=RGBColor(0x40,0x30,0x70))

seq2 = [
    "Prospect submits lead form ──────────▶  createLead()  ──────────▶  INSERT Lead",
    "                                                   ──────────────────────────▶  trigger NexusAILeadScoringFlow",
    "                                                              ───────────────────────────────▶  scoreLead()",
    "                                                              Evaluates: company size, email domain, intent keywords...",
    "                                                              ◀───────────────────────────────  AIScore=78 + Reason",
    "                        UPDATE Lead: AIScore__c=78, AIReason__c='Mid-market IT buyer'",
    "◀─────────────────────  Confirmation email sent to prospect + alert to admin",
    "                                                                          Admin opens Lead record → score ring shown",
    "                                                                          Admin clicks Approve → UPDATE AdminDecision__c",
    "                                                                                        ─────────────────────────────▶  convertLeadFull()",
    "                                                                                        CREATE Account + Contact + Opportunity",
    "◀─────────────────────────────────────────────────────────────────────────────────────  Welcome email + portal access link",
]
ly2 = 2.22
for line in seq2:
    add_text(s, line, 0.2, ly2, 13.0, 0.27, size=8.5, color=RGBColor(0xcd,0xd6,0xf4))
    ly2 += 0.35

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — SEQUENCE: STOCK ALERT
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 14, "Sequence Diagram — Low Stock Alert", "Event-driven auto reorder and dashboard alerting")
add_rect(s, 0, 0.9, 13.33, 6.48, fill=RGBColor(0x1e,0x1e,0x2e))

steps3 = ["1. Stock Drops Low","2. Platform Event Fired","3. Flow Responds","4. PO Auto-Created","5. Notification Sent","6. Dashboard Alerts"]
sx3 = 0.2
for step in steps3:
    add_rect(s, sx3, 1.05, 2.05, 0.38, fill=SF_RED)
    add_text(s, step, sx3+0.05, 1.1, 1.95, 0.28, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    sx3 += 2.18

actors3 = ["External\nSystem", "Stock__c\nRecord", "Platform\nEvent", "LowStock\nFlow", "Purchase\nOrder", "Notif\nSystem", "Dashboard"]
ax3 = [0.15, 1.6, 3.05, 4.6, 6.1, 7.65, 9.2]
for actor, x in zip(actors3, ax3):
    add_rect(s, x, 1.58, 1.3, 0.48, fill=SF_DARK, line=SF_RED, line_w=Pt(1))
    add_text(s, actor, x, 1.6, 1.3, 0.44, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x+0.63, 2.08, 0.02, 4.9, fill=RGBColor(0x60,0x20,0x30))

seq3 = [
    "External system updates stock ──────▶  Qty < Reorder_Threshold__c",
    "               ─────────────────────────────────▶  publish Low_Stock_Alert__e  (Platform Event)",
    "                                    ─────────────────────────────────────────▶  LowStockResponse Flow triggers",
    "                                                  ──────────────────────────────────────────────▶  CREATE Purchase_Order__c",
    "                                                     Status=Pending, Expected_Delivery = Today+5",
    "                                                  ──────────────────────────────────────────────────────────────▶  CREATE Notification",
    "                                                     Type=Warning: 'Pipeline at risk — $24K in active quotes'",
    "",
    "                                                                                                   Admin opens Analytics Dashboard",
    "                                                                                                   getStockHealth() → {critical:1, atRisk:2}",
    "                                                                                                   ◀─────────────────────  Red alert banner shown",
    "",
    "Admin opens Stock War Room → getStockRiskProducts() → AI_Risk_Summary__c shown per product",
]
ly3 = 2.18
for line in seq3:
    add_text(s, line, 0.2, ly3, 13.0, 0.27, size=8.5, color=RGBColor(0xcd,0xd6,0xf4))
    ly3 += 0.33

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — MARKET COMPARISON
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 15, "Market Comparison", "Nexus vs leading B2B platforms")
add_rect(s, 0, 0.9, 13.33, 6.48, fill=LIGHT_BG)

# comparison table
headers = ["Feature", "Nexus ✓", "SF CPQ", "SAP Commerce", "Magento B2B", "HubSpot"]
col_w   = [2.3, 2.2, 1.7, 1.8, 1.8, 1.5]
rows = [
    ["CRM Integration",    "Native (built on SF)", "Native",    "Connector needed", "Plugin",      "Native"],
    ["B2B Quoting",        "Custom builder",        "Full CPQ",  "Via SAP CRM",      "Limited",     "Manual"],
    ["AI Lead Scoring",    "Built-in (rules)",      "Einstein",  "Separate SAP AI",  "None",        "HS Score"],
    ["Stock Intelligence", "War Room + AI Risk",    "Module",    "SAP EWM",          "Basic",       "None"],
    ["E-Signature",        "DocuSign + Dvelop",     "DocuSign",  "Varies",           "Plugins",     "Plugins"],
    ["Analytics",          "Custom KPI+funnel",     "SF Analyt.","SAP Analytics",    "Basic",       "HS Reports"],
    ["Price",              "Dev org (free)",         "$75/user",  "Enterprise",       "$22k+/yr",    "$800+/mo"],
    ["Customizability",    "Very High (LWC+Apex)",  "Medium",    "High (ABAP)",      "High (PHP)",  "Medium"],
]
tx = 0.2
ty = 1.1
# header row
for i, (h, w) in enumerate(zip(headers, col_w)):
    bg = SF_DARK if i == 1 else RGBColor(0x22,0x45,0x60)
    add_rect(s, tx, ty, w, 0.4, fill=bg)
    add_text(s, h, tx+0.05, ty+0.06, w-0.08, 0.3, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx += w + 0.02
ty += 0.42
for ri, row in enumerate(rows):
    tx = 0.2
    bg_row = WHITE if ri % 2 == 0 else RGBColor(0xf5,0xf8,0xfa)
    for i, (cell, w) in enumerate(zip(row, col_w)):
        add_rect(s, tx, ty, w, 0.37, fill=RGBColor(0xe8,0xf4,0xed) if i==1 else bg_row, line=BORDER, line_w=Pt(1))
        color = SF_GREEN if i == 1 else MUTED
        bold = i == 1
        add_text(s, cell, tx+0.05, ty+0.04, w-0.08, 0.29, size=9.5, color=color, bold=bold, align=PP_ALIGN.CENTER)
        tx += w + 0.02
    ty += 0.39

# SWOT
add_text(s, "SWOT", 0.2, 5.15, 1.0, 0.3, size=13, bold=True, color=SF_DARK)
swot = [
    ("Strengths", SF_BLUE, "Native SF · AI everywhere · Dual e-sign · Stock War Room"),
    ("Weaknesses", SF_RED, "Dev org limits · No mobile · Manual stock sync"),
    ("Opportunities", SF_GREEN, "Agentforce · Arabic UI · ERP integration · Mobile PWA"),
    ("Threats", SF_ORANGE, "CPQ overlap · Cost at scale · Governor limits"),
]
sx4 = 0.2
for title, color, desc in swot:
    add_rect(s, sx4, 5.5, 3.1, 0.9, fill=WHITE, line=color, line_w=Pt(2))
    add_text(s, title, sx4+0.1, 5.54, 2.9, 0.28, size=11, bold=True, color=color)
    add_text(s, desc, sx4+0.1, 5.82, 2.9, 0.5, size=9.5, color=MUTED)
    sx4 += 3.27

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, 16, "Conclusion", "What Nexus delivers and what comes next")
add_rect(s, 0, 0.9, 13.33, 6.48, fill=LIGHT_BG)

add_rect(s, 0.3, 1.05, 12.73, 0.75, fill=SF_DARK)
add_text(s, '"Nexus proves that Salesforce is not just a CRM — it is a complete platform for enterprise-grade B2B commerce,'
            ' AI-driven sales intelligence, and end-to-end contract lifecycle management."',
         0.45, 1.1, 12.5, 0.65, size=11.5, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

add_bullet_box(s, [
    "60+ LWC Components — full UI from scratch",
    "30+ Apex Classes — controllers + service layer",
    "6 Custom Objects + 50+ custom fields",
    "9 Automation Flows (lead scoring, stock alerts, order tracking)",
    "AI at every touchpoint: leads, quotes, stock, recommendations",
    "4 External integrations: DocuSign, Dvelop, Stripe, Agentforce",
    "Zero external backend — 100% Salesforce stack",
], 0.3, 1.95, 6.1, 3.1, title="What Was Delivered", accent=SF_BLUE, title_color=SF_BLUE)

add_bullet_box(s, [
    "Governor Limits demand bulkification from day one",
    "@wire decorators are powerful but need cache strategy",
    "Flows + Apex is the sweet spot for automation",
    "Custom Metadata Types beat hardcoded configs every time",
    "Portal FLS differs from internal — test both separately",
], 6.83, 1.95, 6.2, 2.1, title="Lessons Learned", accent=SF_GREEN, title_color=SF_GREEN)

add_text(s, "Roadmap", 6.83, 4.18, 3, 0.3, size=13, bold=True, color=SF_DARK)
roadmap = [
    ("Short-term  (1–3 mo)", "Mobile UI · Arabic language · ERP integration"),
    ("Medium-term (3–6 mo)", "ML demand forecasting · Multi-currency · Supplier portal"),
    ("Long-term   (6–12 mo)","Multi-org enterprise · Marketplace · Blockchain audit"),
]
ry = 4.55
for period, items in roadmap:
    add_rect(s, 6.83, ry, 6.2, 0.55, fill=WHITE, line=BORDER, line_w=Pt(1))
    add_text(s, period, 6.93, ry+0.04, 2.2, 0.26, size=10, bold=True, color=SF_DARK)
    add_text(s, items,  9.15, ry+0.04, 3.8, 0.26, size=10, color=MUTED)
    ry += 0.58

add_text(s, "Stack: LWC · Apex · Flows · Einstein AI · DocuSign · Dvelop · Stripe · Agentforce",
         0.3, 7.1, 12.73, 0.25, size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — THANK YOU
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_DARK)
add_rect(s, 0, 0, 6.67, 7.5, fill=RGBColor(0x0b,0x53,0x94))

tb2 = s.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.33), Inches(1.4))
tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r = p2.add_run(); r.text = "Thank You"
r.font.size = Pt(64); r.font.bold = True; r.font.color.rgb = WHITE

add_text(s, "Nexus — B2B E-Commerce Platform on Salesforce",
         1, 3.6, 11.33, 0.45, size=18, color=RGBColor(0xb0,0xc8,0xe8), align=PP_ALIGN.CENTER)
add_text(s, "Built with  LWC · Apex · Flows · Einstein AI · DocuSign · Dvelop · Stripe",
         1, 4.15, 11.33, 0.35, size=14, color=RGBColor(0x80,0xa8,0xd0), align=PP_ALIGN.CENTER)
add_text(s, "March 2026",
         1, 5.0, 11.33, 0.35, size=13, color=RGBColor(0x60,0x80,0xa0), align=PP_ALIGN.CENTER)

# ─── Save ──────────────────────────────────────────────────────────────────
out = r"c:\Users\MSI\Desktop\Salesforces\ecommerce-platform\PRESENTATION_NEXUS.pptx"
prs.save(out)
print(f"Saved: {out}")
